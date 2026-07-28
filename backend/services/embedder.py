import datetime
import io
import re
import uuid
from typing import Any, Dict, List, Optional

import pypdf
from services.settings import settings

# Multi-format document parser imports with graceful fallback
try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    SentenceTransformer = None

# Retained for the legacy pgvector path. P2 will replace it with the frozen
# Qwen3-Embedding-0.6B model behind the same service interface.
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"


class DocumentProcessor:
    def __init__(self, chunk_size: Optional[int] = None, overlap: Optional[int] = None):
        self.chunk_size = chunk_size if chunk_size is not None else settings.chunk_size_chars
        self.overlap = overlap if overlap is not None else settings.chunk_overlap_chars

    def extract_chunks(
        self,
        file_bytes: bytes,
        filename: str,
        category: Optional[str] = None,
        doc_id: Optional[str] = None
    ) -> List[Dict[str, Any]]:
        """
        Extracts text from PDF, DOCX, XLSX, PPTX, or TXT files, chunks text,
        and returns metadata-tagged chunk dictionaries.
        """
        if not doc_id:
            doc_id = str(uuid.uuid4())

        ext = filename.lower().split('.')[-1]
        pages_text = []

        if ext == 'pdf':
            pages_text = self._parse_pdf(file_bytes)
        elif ext == 'docx':
            pages_text = self._parse_docx(file_bytes)
        elif ext == 'pptx':
            pages_text = self._parse_pptx(file_bytes)
        elif ext in ['xlsx', 'xls']:
            pages_text = self._parse_xlsx(file_bytes)
        else:
            # Fallback for plain text / markdown files
            raw_text = file_bytes.decode('utf-8', errors='ignore')
            pages_text = [raw_text]

        num_pages = len(pages_text) or 1
        chunks = []
        global_chunk_idx = 0

        for page_idx, raw_text in enumerate(pages_text):
            page_num = page_idx + 1
            cleaned_text = self._clean_text(raw_text)

            if not cleaned_text:
                continue

            page_chunks = self._chunk_text(cleaned_text)
            for idx_in_page, chunk_text in enumerate(page_chunks):
                chunks.append({
                    "id": f"{doc_id}_p{page_num}_c{idx_in_page}",
                    "doc_id": doc_id,
                    "filename": filename,
                    "category": category or settings.default_category,
                    "page_number": page_num,
                    "total_pages": num_pages,
                    "chunk_index": global_chunk_idx,
                    "text": chunk_text,
                    "created_at": datetime.datetime.utcnow().isoformat()
                })
                global_chunk_idx += 1

        return chunks

    def _parse_pdf(self, file_bytes: bytes) -> List[str]:
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        return [(p.extract_text() or "") for p in reader.pages]

    def _parse_docx(self, file_bytes: bytes) -> List[str]:
        if not docx:
            return [""]
        doc = docx.Document(io.BytesIO(file_bytes))
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        return [full_text]

    def _parse_pptx(self, file_bytes: bytes) -> List[str]:
        if not pptx:
            return [""]
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        slides_text = []
        for slide in prs.slides:
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text.strip())
            slides_text.append("\n".join(slide_content))
        return slides_text

    def _parse_xlsx(self, file_bytes: bytes) -> List[str]:
        if not openpyxl:
            return [""]
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        sheets_text = []
        for sheet in wb.worksheets:
            rows_str = []
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(cell) for cell in row if cell is not None]
                if row_vals:
                    rows_str.append(" | ".join(row_vals))
            sheets_text.append(f"Sheet: {sheet.title}\n" + "\n".join(rows_str))
        return sheets_text

    def _clean_text(self, text: str) -> str:
        return re.sub(r"\s+", " ", text).strip()

    def _chunk_text(self, text: str) -> List[str]:
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + self.chunk_size
            chunk = text[start:end]

            if end < text_len:
                last_space = chunk.rfind(" ")
                if last_space > self.chunk_size // 2:
                    end = start + last_space
                    chunk = text[start:end]

            chunks.append(chunk.strip())
            start = end - self.overlap if (end - self.overlap) > start else end

        return chunks


# Alias for backward compatibility
PDFProcessor = DocumentProcessor


class EmbeddingService:
    def __init__(
        self,
        model: Optional[Any] = None,
        model_name: str = EMBEDDING_MODEL_NAME,
        embedding_dimension: int = 384,
    ):
        """Create an embedding service without network/model-load side effects.

        Loading a Hugging Face model at module import made offline startup hang
        while waiting for a cache or network response. The production model
        will be loaded explicitly during P2 index startup; tests and the
        current legacy path use the deterministic fallback until then.
        """
        self.model = model
        self.model_name = model_name
        self.embedding_dimension = embedding_dimension

    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        if self.model is not None:
            embeddings = self.model.encode(texts, show_progress_bar=False, convert_to_numpy=True)
            return embeddings.tolist()
        else:
            return [self._fallback_vector(t) for t in texts]

    def load_local_model(self, model_path: str, revision: Optional[str] = None) -> None:
        """Load a cached SentenceTransformers model without network access."""
        if SentenceTransformer is None:
            raise RuntimeError("sentence-transformers is not installed")
        self.model = SentenceTransformer(
            model_path,
            device="cpu",
            local_files_only=True,
            trust_remote_code=True,
            revision=revision,
        )
        self.embedding_dimension = self.model.get_sentence_embedding_dimension()

    def _fallback_vector(self, text: str, dim: Optional[int] = None) -> List[float]:
        import hashlib
        import math

        dim = dim or self.embedding_dimension
        vec = []
        for i in range(dim):
            h = hashlib.sha256(f"{text}_{i}".encode("utf-8")).hexdigest()
            val = (int(h[:8], 16) / 0xFFFFFFFF) * 2.0 - 1.0
            vec.append(val)
        norm = math.sqrt(sum(x * x for x in vec)) or 1.0
        return [x / norm for x in vec]


class QwenEmbeddingService(EmbeddingService):
    """Qwen3-Embedding adapter used by the frozen P2 dense index."""

    MODEL_NAME = "Qwen/Qwen3-Embedding-0.6B"

    def __init__(
        self,
        model: Optional[Any] = None,
        model_name: str = MODEL_NAME,
        embedding_dimension: int = 1024,
    ):
        super().__init__(
            model=model,
            model_name=model_name,
            embedding_dimension=embedding_dimension,
        )
        self.model_revision: Optional[str] = None

    @property
    def is_ready(self) -> bool:
        return self.model is not None

    def load_local_model(self, model_path: str, revision: Optional[str] = None) -> None:
        super().load_local_model(model_path, revision=revision)
        self.model_revision = revision or model_path

    def generate_embeddings(self, texts: List[str]) -> List[List[float]]:
        if self.model is None:
            return [self._fallback_vector(text) for text in texts]

        try:
            embeddings = self.model.encode(
                texts,
                batch_size=8,
                show_progress_bar=False,
                convert_to_numpy=True,
                normalize_embeddings=True,
            )
        except TypeError:
            embeddings = self.model.encode(
                texts,
                batch_size=8,
                show_progress_bar=False,
                convert_to_numpy=True,
            )
        rows = embeddings.tolist()
        normalized = []
        for row in rows:
            norm = sum(value * value for value in row) ** 0.5 or 1.0
            normalized.append([value / norm for value in row])
        return normalized
